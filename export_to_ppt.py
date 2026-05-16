import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from bs4 import BeautifulSoup

BASE_DIR = Path(__file__).resolve().parent
HTML_PATH = BASE_DIR / "index.html"
OUTPUT_PATH = BASE_DIR / "output.pptx"

# Simple helper to add a title+bullets slide
def add_bullet_slide(prs, title, bullets):
    if not bullets:
        return
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
    slide.shapes.title.text = title
    body = slide.shapes.placeholders[1].text_frame
    body.text = bullets[0]
    for b in bullets[1:]:
        body.add_paragraph().text = b

# Helper to add an image slide with caption

def add_image_slide(prs, title, caption, image_path):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    left = Inches(0.5)
    top = Inches(0.3)
    width = prs.slide_width - Inches(1)
    # Title
    title_box = slide.shapes.add_textbox(left, top, width, Inches(0.6))
    title_tf = title_box.text_frame
    title_tf.text = title
    title_tf.paragraphs[0].font.size = Pt(24)
    current_top = top + Inches(0.8)

    if image_path and image_path.exists():
        max_width = prs.slide_width - Inches(1)
        pic = slide.shapes.add_picture(str(image_path), left, current_top, width=max_width)
        current_top = current_top + pic.height + Inches(0.2)
    else:
        current_top = current_top + Inches(0.2)

    if caption:
        cap_box = slide.shapes.add_textbox(left, current_top, width, Inches(1))
        cap_box.text = caption


def collect_start_cards(soup):
    bullets = []
    for card in soup.select('.start-card'):
        title = card.select_one('.card-title').get_text(strip=True)
        descs = [d.get_text(strip=True).lstrip('-').strip() for d in card.select('.card-desc')]
        text = f"{title}: " + "; ".join(descs)
        bullets.append(text)
    return bullets


def collect_journey_cards(container):
    bullets = []
    for card in container.select('.journey-card'):
        title = card.select_one('.card-title').get_text(strip=True)
        desc = card.select_one('p').get_text(strip=True)
        bullets.append(f"{title}: {desc}")
    return bullets


def collect_reflection_cards(container):
    bullets = []
    for card in container.select('.reflection-card'):
        title = card.select_one('.card-title').get_text(strip=True)
        desc = card.select_one('p').get_text(strip=True)
        bullets.append(f"{title}: {desc}")
    return bullets


def add_section_slides(prs, section, section_title):
    # Tab contents
    for tab in section.select('.tab-content'):
        tab_id = tab.get('id', '')
        tab_label = tab_id.split('-')[-1] if '-' in tab_id else ''
        label_map = {
            'challenges': 'Pain Points',
            'results': 'Solutions',
            'challenge': 'Biggest Challenge',
            'reflection': 'Key Reflection'
        }
        sub_title = label_map.get(tab_label, tab_label.title() if tab_label else '')

        # Carousel items with images
        for item in tab.select('.carousel-item'):
            caption_el = item.select_one('.image-caption')
            caption = caption_el.get_text(strip=True) if caption_el else ''
            img_el = item.select_one('img')
            img_src = img_el['src'] if img_el and img_el.has_attr('src') else ''
            img_path = BASE_DIR / img_src if img_src else None
            title = f"{section_title} — {sub_title}" if sub_title else section_title
            add_image_slide(prs, title, caption, img_path)

        # Standalone images (non-carousel)
        for img_el in tab.select('.station-images img'):
            img_src = img_el['src'] if img_el.has_attr('src') else ''
            img_path = BASE_DIR / img_src if img_src else None
            caption = img_el.get('alt', '')
            title = f"{section_title} — {sub_title}" if sub_title else section_title
            add_image_slide(prs, title, caption, img_path)

        # Paragraph-only content
        paragraphs = [p.get_text(strip=True) for p in tab.find_all('p') if not p.find_parent(class_='carousel-item')]
        if paragraphs:
            title = f"{section_title} — {sub_title}" if sub_title else section_title
            add_bullet_slide(prs, title, paragraphs)

        # Reflection cards inside this tab
        reflection_cards = tab.select('.reflection-cards')
        for rc in reflection_cards:
            bullets = collect_reflection_cards(rc)
            if bullets:
                title = f"{section_title} — {sub_title}" if sub_title else section_title
                add_bullet_slide(prs, title, bullets)


def main():
    html = HTML_PATH.read_text(encoding='utf-8')
    soup = BeautifulSoup(html, 'lxml')

    prs = Presentation()

    # Title slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = "My AI Journey"
    subtitle = title_slide.placeholders[1]
    subtitle.text = "LBS MBA Interview Presentation"

    # Start cards
    start_bullets = collect_start_cards(soup)
    add_bullet_slide(prs, "Start — Foundations", start_bullets)

    # Journey sections
    for sec in soup.select('section.journey-section'):
        sec_id = sec.get('id')
        sec_title_el = sec.find('h3')
        sec_title = sec_title_el.get_text(strip=True) if sec_title_el else 'Section'

        if sec_id in {"station1", "station2", "station3"}:
            add_section_slides(prs, sec, sec_title)
        elif sec_id == 'lookout':
            cards = sec.select('.lookout-card')
            bullets = []
            for card in cards:
                title = card.select_one('.card-title').get_text(strip=True)
                desc = card.select_one('p').get_text(strip=True)
                bullets.append(f"{title}: {desc}")
            add_bullet_slide(prs, sec_title, bullets)
        elif sec_id == 'new-journey':
            # Journey cards under next station
            cards_container = sec.select_one('.next-station-column')
            if cards_container:
                bullets = collect_journey_cards(cards_container)
                add_bullet_slide(prs, sec_title, bullets)
        elif sec_id == 'summary':
            summary_text = [sec_title]
            add_bullet_slide(prs, "Summary", summary_text)

    prs.save(OUTPUT_PATH)
    print(f"Saved to {OUTPUT_PATH}")


if __name__ == "__main__":
    main()
